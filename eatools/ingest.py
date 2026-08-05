"""Turn an uploaded file into a normalised :class:`SourceDoc`.

`ingest(filename, data)` dispatches on the file suffix. **Adding a new input format
is a single function here plus one line in `_PARSERS` — nothing downstream changes**,
because every parser returns the same `SourceDoc` shape.

Structured formats (draw.io, Visio, PowerPoint) become a graph of shapes + connectors
that `SourceDoc.render()` flattens to text. Page-image formats (PDF, images) become
base64 PNGs for a vision model. Anything unrecognised — or any parser that chokes on a
malformed file — raises :class:`UnsupportedFile`, which `app.py` catches so one bad file
never kills a multi-file batch.
"""

from __future__ import annotations

import base64
import io
import re
import zipfile
import zlib
from urllib.parse import unquote
from xml.etree import ElementTree as ET

from .model import Connector, Page, Shape, SourceDoc

# Long edge (px) images are clamped to before sending to the model. Matches the
# high-resolution vision ceiling of the current Claude models.
MAX_IMAGE_EDGE = 2576
PDF_DPI = 150
PDF_PAGE_CAP = 20

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".tif"}


class UnsupportedFile(Exception):
    """Raised for an unknown suffix or a file that cannot be parsed."""


# ---------------------------------------------------------------------------
# draw.io / mxGraph
# ---------------------------------------------------------------------------

_HTML_TAG = re.compile(r"<[^>]+>")
_WS = re.compile(r"\s+")


def _strip_html(text: str) -> str:
    """mxGraph labels are HTML; reduce to readable plain text."""
    if not text:
        return ""
    text = text.replace("<br>", " ").replace("<br/>", " ").replace("&nbsp;", " ")
    text = _HTML_TAG.sub("", text)
    # Unescape the handful of entities mxGraph commonly emits.
    for entity, char in (("&amp;", "&"), ("&lt;", "<"), ("&gt;", ">"), ("&quot;", '"'), ("&#39;", "'")):
        text = text.replace(entity, char)
    return _WS.sub(" ", text).strip()


def _kind_from_style(style: str) -> str:
    """Derive a coarse shape kind from an mxGraph style string."""
    if not style:
        return "shape"
    s = style.lower()
    if "swimlane" in s:
        return "swimlane"
    if "shape=cylinder" in s or "cylinder" in s:
        return "database"
    if "shape=cloud" in s or "cloud" in s:
        return "cloud"
    if "ellipse" in s:
        return "ellipse"
    if "rhombus" in s:
        return "decision"
    if "container=1" in s or "group" in s:
        return "container"
    return "shape"


def _decode_drawio_payload(text: str) -> str:
    """draw.io stores diagrams as deflate+base64+URL-encoded text inside <diagram>.

    Returns the inner mxGraphModel XML. Falls back to the raw text when it is already
    plain XML (uncompressed export).
    """
    stripped = text.strip()
    if stripped.startswith("<mxGraphModel"):
        return stripped
    try:
        raw = base64.b64decode(stripped)
        # raw DEFLATE (no zlib header) — negative wbits.
        inflated = zlib.decompress(raw, -15)
        return unquote(inflated.decode("utf-8"))
    except Exception:
        # Not a compressed payload; hand back what we were given.
        return stripped


def _parse_mxgraph(model_xml: str, page_name: str) -> Page:
    root = ET.fromstring(model_xml)
    # <mxGraphModel><root><mxCell .../></root></mxGraphModel>
    cells = root.iter("mxCell")
    shapes: list[Shape] = []
    edges: list[tuple[str, str | None, str | None, str]] = []  # label, source, target, id

    default_parents = set()
    for cell in root.iter("mxCell"):
        cid = cell.get("id", "")
        if cid in ("0", "1"):
            default_parents.add(cid)

    for cell in cells:
        cid = cell.get("id", "")
        if cid in ("0", "1"):
            continue
        style = cell.get("style", "") or ""
        label = _strip_html(cell.get("value", "") or "")

        if cell.get("edge") == "1":
            edges.append((label, cell.get("source"), cell.get("target"), cid))
            continue

        if cell.get("vertex") == "1" or label:
            parent = cell.get("parent")
            if parent in default_parents:
                parent = None
            geometry = None
            geo = cell.find("mxGeometry")
            if geo is not None:
                try:
                    geometry = (
                        float(geo.get("x", 0)),
                        float(geo.get("y", 0)),
                        float(geo.get("width", 0)),
                        float(geo.get("height", 0)),
                    )
                except (TypeError, ValueError):
                    geometry = None
            shapes.append(
                Shape(id=cid, text=label, kind=_kind_from_style(style), parent=parent, geometry=geometry)
            )

    known_ids = {s.id for s in shapes}
    connectors = [
        Connector(source=src, target=tgt, label=lbl)
        for (lbl, src, tgt, _cid) in edges
        # keep an edge even if one endpoint is unknown; render() copes with None
        if src in known_ids or tgt in known_ids or lbl
    ]
    # Drop parent links that point at an edge or a missing shape.
    for shape in shapes:
        if shape.parent and shape.parent not in known_ids:
            shape.parent = None
    return Page(shapes=shapes, connectors=connectors, name=page_name)


def _ingest_drawio(filename: str, data: bytes) -> SourceDoc:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise UnsupportedFile(f"{filename}: not valid UTF-8 XML") from exc
    try:
        outer = ET.fromstring(text)
    except ET.ParseError as exc:
        raise UnsupportedFile(f"{filename}: malformed XML ({exc})") from exc

    pages: list[Page] = []
    diagrams = list(outer.iter("diagram"))
    if diagrams:
        for i, diagram in enumerate(diagrams, 1):
            name = diagram.get("name") or str(i)
            # Uncompressed exports nest <mxGraphModel> as a child element; compressed
            # ones store the deflate+base64 payload as <diagram> text.
            model = diagram.find("mxGraphModel")
            if model is not None:
                inner = ET.tostring(model, encoding="unicode")
            else:
                inner = _decode_drawio_payload(diagram.text or "")
            try:
                pages.append(_parse_mxgraph(inner, name))
            except ET.ParseError as exc:
                raise UnsupportedFile(f"{filename}: unreadable diagram '{name}' ({exc})") from exc
    elif outer.tag == "mxGraphModel":
        pages.append(_parse_mxgraph(text, "1"))
    else:
        # A bare .xml that happens to hold mxGraphModel somewhere.
        model = outer.find(".//mxGraphModel")
        if model is None:
            raise UnsupportedFile(f"{filename}: no mxGraphModel found")
        pages.append(_parse_mxgraph(ET.tostring(model, encoding="unicode"), "1"))

    if not any(p.shapes for p in pages):
        raise UnsupportedFile(f"{filename}: diagram contains no shapes")
    return SourceDoc(kind="structured", source_name=filename, pages=pages)


# ---------------------------------------------------------------------------
# Visio (.vsdx)
# ---------------------------------------------------------------------------

_VISIO_NS = {"v": "http://schemas.microsoft.com/office/visio/2012/main"}


def _ingest_vsdx(filename: str, data: bytes) -> SourceDoc:
    try:
        zf = zipfile.ZipFile(io.BytesIO(data))
    except zipfile.BadZipFile as exc:
        raise UnsupportedFile(f"{filename}: not a valid .vsdx (zip) file") from exc

    page_files = sorted(
        n for n in zf.namelist() if re.match(r"visio/pages/page\d+\.xml$", n)
    )
    if not page_files:
        raise UnsupportedFile(f"{filename}: no visio/pages/*.xml found")

    pages: list[Page] = []
    for idx, pf in enumerate(page_files, 1):
        try:
            root = ET.fromstring(zf.read(pf))
        except ET.ParseError as exc:
            raise UnsupportedFile(f"{filename}: unreadable page {pf} ({exc})") from exc

        shapes_by_id: dict[str, Shape] = {}
        for shape_el in root.iter("{%s}Shape" % _VISIO_NS["v"]):
            sid = shape_el.get("ID")
            if sid is None:
                continue
            text_el = shape_el.find("v:Text", _VISIO_NS)
            text = "".join(text_el.itertext()).strip() if text_el is not None else ""
            shapes_by_id[sid] = Shape(id=sid, text=text, kind="shape")

        # <Connect FromSheet=".." ToSheet=".." FromCell="BeginX|EndX"/>
        # FromSheet is the connector shape; ToSheet an endpoint. BeginX = source end,
        # EndX = target end. Group the two Connect rows that share a connector shape.
        connector_ends: dict[str, dict[str, str]] = {}
        for conn in root.iter("{%s}Connect" % _VISIO_NS["v"]):
            from_sheet = conn.get("FromSheet")
            to_sheet = conn.get("ToSheet")
            from_cell = conn.get("FromCell", "")
            if not from_sheet or not to_sheet:
                continue
            ends = connector_ends.setdefault(from_sheet, {})
            if from_cell.startswith("Begin"):
                ends["source"] = to_sheet
            elif from_cell.startswith("End"):
                ends["target"] = to_sheet
            else:
                ends.setdefault("source" if "source" not in ends else "target", to_sheet)

        connectors: list[Connector] = []
        for conn_shape_id, ends in connector_ends.items():
            label = ""
            conn_shape = shapes_by_id.pop(conn_shape_id, None)  # drop connectors from nodes
            if conn_shape is not None:
                label = conn_shape.text
            connectors.append(
                Connector(source=ends.get("source"), target=ends.get("target"), label=label)
            )

        pages.append(
            Page(shapes=list(shapes_by_id.values()), connectors=connectors, name=str(idx))
        )

    if not any(p.shapes for p in pages):
        raise UnsupportedFile(f"{filename}: no shapes with text found")
    return SourceDoc(kind="structured", source_name=filename, pages=pages)


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------


def _ingest_pptx(filename: str, data: bytes) -> SourceDoc:
    try:
        from pptx import Presentation  # noqa: PLC0415 (optional heavy import)
        from pptx.util import Emu  # noqa: F401,PLC0415
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedFile("python-pptx is not installed") from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise UnsupportedFile(f"{filename}: not a valid .pptx ({exc})") from exc

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    pages: list[Page] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        shapes: list[Shape] = []
        connectors: list[Connector] = []
        for shape in slide.shapes:
            sid = str(shape.shape_id)
            text = (shape.text_frame.text.strip() if shape.has_text_frame else "") if hasattr(shape, "has_text_frame") else ""

            # Connector shapes carry <a:stCxn>/<a:endCxn> naming the shapes they join.
            spxml = shape._element  # noqa: SLF001 (python-pptx has no public accessor)
            st = spxml.find(".//a:stCxn", ns)
            end = spxml.find(".//a:endCxn", ns)
            if st is not None or end is not None:
                connectors.append(
                    Connector(
                        source=st.get("id") if st is not None else None,
                        target=end.get("id") if end is not None else None,
                        label=text,
                    )
                )
                continue

            if text:
                shapes.append(Shape(id=sid, text=text, kind="shape"))

        # Speaker notes carry architect intent; attach as a pseudo-shape.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                shapes.append(Shape(id=f"notes-{slide_idx}", text=f"[Notes] {notes}", kind="notes"))

        pages.append(Page(shapes=shapes, connectors=connectors, name=str(slide_idx)))

    if not any(p.shapes for p in pages):
        raise UnsupportedFile(f"{filename}: no text shapes found")
    return SourceDoc(kind="structured", source_name=filename, pages=pages)


# ---------------------------------------------------------------------------
# PDF and images → base64 PNG (visual)
# ---------------------------------------------------------------------------


def _ingest_pdf(filename: str, data: bytes) -> SourceDoc:
    try:
        import fitz  # PyMuPDF  # noqa: PLC0415
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedFile("PyMuPDF is not installed") from exc

    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:
        raise UnsupportedFile(f"{filename}: not a valid PDF ({exc})") from exc

    zoom = PDF_DPI / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    images: list[str] = []
    for page in doc[:PDF_PAGE_CAP]:
        pix = page.get_pixmap(matrix=matrix)
        images.append(base64.b64encode(pix.tobytes("png")).decode("ascii"))
    doc.close()

    if not images:
        raise UnsupportedFile(f"{filename}: PDF has no pages")
    return SourceDoc(kind="visual", source_name=filename, images=images)


def _ingest_image(filename: str, data: bytes) -> SourceDoc:
    try:
        from PIL import Image  # noqa: PLC0415
    except ImportError as exc:  # pragma: no cover
        raise UnsupportedFile("Pillow is not installed") from exc

    try:
        img = Image.open(io.BytesIO(data))
        img.load()
    except Exception as exc:
        raise UnsupportedFile(f"{filename}: not a readable image ({exc})") from exc

    if img.mode not in ("RGB", "RGBA", "L"):
        img = img.convert("RGB")
    long_edge = max(img.size)
    if long_edge > MAX_IMAGE_EDGE:
        scale = MAX_IMAGE_EDGE / long_edge
        img = img.resize((max(1, int(img.width * scale)), max(1, int(img.height * scale))))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    return SourceDoc(kind="visual", source_name=filename, images=[b64])


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

_PARSERS = {
    ".drawio": _ingest_drawio,
    ".xml": _ingest_drawio,
    ".vsdx": _ingest_vsdx,
    ".pptx": _ingest_pptx,
    ".pdf": _ingest_pdf,
}


def _suffix(filename: str) -> str:
    dot = filename.rfind(".")
    return filename[dot:].lower() if dot != -1 else ""


def ingest(filename: str, data: bytes) -> SourceDoc:
    """Normalise one uploaded file into a SourceDoc, or raise UnsupportedFile."""
    if not data:
        raise UnsupportedFile(f"{filename}: empty file")
    suffix = _suffix(filename)
    if suffix in _PARSERS:
        return _PARSERS[suffix](filename, data)
    if suffix in IMAGE_SUFFIXES:
        return _ingest_image(filename, data)
    raise UnsupportedFile(f"{filename}: unsupported file type '{suffix or 'unknown'}'")
