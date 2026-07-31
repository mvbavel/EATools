"""Verify each ingester against a programmatically-generated sample of its format."""

import base64
import io
import sys
import zipfile
import zlib
from pathlib import Path
from urllib.parse import quote

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from eatools.ingest import UnsupportedFile, ingest  # noqa: E402


def _check(name, doc, kind):
    assert doc.kind == kind, f"{name}: expected {kind}, got {doc.kind}"
    print(f"  OK  {name}: kind={doc.kind}", end="")
    if kind == "structured":
        shapes = sum(len(p.shapes) for p in doc.pages)
        conns = sum(len(p.connectors) for p in doc.pages)
        print(f" pages={len(doc.pages)} shapes={shapes} connectors={conns}")
        rendered = doc.render()
        assert rendered, f"{name}: empty render"
    else:
        print(f" images={len(doc.images)}")
        assert doc.images, f"{name}: no images"


# --- draw.io inline ---------------------------------------------------------
DRAWIO_INLINE = """<mxfile><diagram name="Page-1"><mxGraphModel><root>
<mxCell id="0"/><mxCell id="1" parent="0"/>
<mxCell id="2" value="&lt;b&gt;Billing&lt;/b&gt;" style="rounded=1" vertex="1" parent="1"><mxGeometry x="10" y="10" width="120" height="60"/></mxCell>
<mxCell id="3" value="PostgreSQL" style="shape=cylinder" vertex="1" parent="1"><mxGeometry x="10" y="120" width="80" height="80"/></mxCell>
<mxCell id="4" value="writes" style="edgeStyle" edge="1" parent="1" source="2" target="3"/>
</root></mxGraphModel></diagram></mxfile>"""


def test_drawio_inline():
    doc = ingest("sample.drawio", DRAWIO_INLINE.encode())
    _check("drawio-inline", doc, "structured")
    shapes = doc.pages[0].shapes
    assert any(s.text == "Billing" for s in shapes), "HTML not stripped"
    assert any(s.kind == "database" for s in shapes), "cylinder kind not derived"
    assert doc.pages[0].connectors[0].label == "writes"


# --- draw.io compressed -----------------------------------------------------
def test_drawio_compressed():
    inner = """<mxGraphModel><root><mxCell id="0"/><mxCell id="1" parent="0"/>
<mxCell id="2" value="CRM" vertex="1" parent="1"><mxGeometry x="0" y="0" width="80" height="40"/></mxCell>
</root></mxGraphModel>"""
    compressor = zlib.compressobj(9, zlib.DEFLATED, -15)
    deflated = compressor.compress(quote(inner).encode()) + compressor.flush()
    payload = base64.b64encode(deflated).decode()
    xml = f'<mxfile><diagram name="P1">{payload}</diagram></mxfile>'
    doc = ingest("compressed.drawio", xml.encode())
    _check("drawio-compressed", doc, "structured")
    assert any(s.text == "CRM" for s in doc.pages[0].shapes)


# --- plain .xml -------------------------------------------------------------
def test_xml():
    doc = ingest("diagram.xml", DRAWIO_INLINE.encode())
    _check("xml", doc, "structured")


# --- vsdx -------------------------------------------------------------------
def _make_vsdx():
    ns = "http://schemas.microsoft.com/office/visio/2012/main"
    page = f"""<?xml version="1.0"?>
<PageContents xmlns="{ns}">
  <Shapes>
    <Shape ID="1"><Text>Order Service</Text></Shape>
    <Shape ID="2"><Text>Inventory</Text></Shape>
    <Shape ID="3"><Text>calls</Text></Shape>
  </Shapes>
  <Connects>
    <Connect FromSheet="3" ToSheet="1" FromCell="BeginX"/>
    <Connect FromSheet="3" ToSheet="2" FromCell="EndX"/>
  </Connects>
</PageContents>"""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("visio/pages/page1.xml", page)
    return buf.getvalue()


def test_vsdx():
    doc = ingest("sample.vsdx", _make_vsdx())
    _check("vsdx", doc, "structured")
    shapes = doc.pages[0].shapes
    texts = {s.text for s in shapes}
    assert texts == {"Order Service", "Inventory"}, f"connector not dropped: {texts}"
    conn = doc.pages[0].connectors[0]
    assert conn.source == "1" and conn.target == "2", (conn.source, conn.target)
    assert conn.label == "calls"


# --- pptx -------------------------------------------------------------------
def _make_pptx():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(0, 0, 100, 50)
    box.text_frame.text = "Payment Gateway"
    box2 = slide.shapes.add_textbox(0, 100, 100, 50)
    box2.text_frame.text = "Fraud Check"
    slide.notes_slide.notes_text_frame.text = "Both are SaaS."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx():
    doc = ingest("sample.pptx", _make_pptx())
    _check("pptx", doc, "structured")
    texts = " ".join(s.text for s in doc.pages[0].shapes)
    assert "Payment Gateway" in texts
    assert "Both are SaaS" in texts, "speaker notes not included"


# --- pdf --------------------------------------------------------------------
def _make_pdf():
    import fitz

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Architecture Overview")
    data = doc.tobytes()
    doc.close()
    return data


def test_pdf():
    doc = ingest("sample.pdf", _make_pdf())
    _check("pdf", doc, "visual")


# --- image ------------------------------------------------------------------
def _make_png():
    from PIL import Image

    img = Image.new("RGB", (4000, 1000), "white")
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def test_image():
    doc = ingest("sample.png", _make_png())
    _check("image", doc, "visual")
    # long edge (4000) should be clamped to <= 2576
    raw = base64.b64decode(doc.images[0])
    from PIL import Image

    w, h = Image.open(io.BytesIO(raw)).size
    assert max(w, h) <= 2576, f"not clamped: {w}x{h}"


# --- bad file ---------------------------------------------------------------
def test_bad_file():
    try:
        ingest("broken.drawio", b"not xml at all <<<")
        raise AssertionError("expected UnsupportedFile")
    except UnsupportedFile as exc:
        print(f"  OK  bad-file raised UnsupportedFile: {exc}")

    try:
        ingest("mystery.zzz", b"data")
        raise AssertionError("expected UnsupportedFile")
    except UnsupportedFile as exc:
        print(f"  OK  unknown-suffix raised UnsupportedFile: {exc}")


if __name__ == "__main__":
    print("Ingest tests:")
    test_drawio_inline()
    test_drawio_compressed()
    test_xml()
    test_vsdx()
    test_pptx()
    test_pdf()
    test_image()
    test_bad_file()
    print("All ingest tests passed.")
